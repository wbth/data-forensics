from flask import Flask, render_template, request, redirect, url_for, session, flash, send_from_directory
from pikepdf._core import PdfError
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader  # Updated to PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.cell.cell import Cell
from PIL import Image, ExifTags, ImageStat
from time import time
from datetime import datetime
from flask_paginate import Pagination, get_page_parameter  # Untuk pagination
from flask_socketio import SocketIO, emit
from pdf2image import convert_from_path
from sklearn.ensemble import IsolationForest
from sklearn.preprocessing import StandardScaler
import os
import hashlib
import json
import joblib
import shutil
import logging
import requests
import pikepdf
import pytesseract
import mimetypes
import csv
import numpy as np
import pandas as pd


# Configure logging to save error log inside the logs folder
logging.basicConfig(filename='logs/app_error.log', level=logging.ERROR)

app = Flask(__name__)
app.secret_key = 'your_secret_key'
#UPLOAD_FOLDER = 'uploads/'
PROCESSED_FOLDER = 'processed/'
DATASET_FOLDER = 'dataset/'

MODEL_PATH = os.path.join(DATASET_FOLDER, 'anomaly_model.pkl')
SCALER_PATH = os.path.join(DATASET_FOLDER, 'scaler.pkl')
ANOMALY_LOG_PATH = os.path.join(DATASET_FOLDER, 'anomaly_log.csv')

#app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['PROCESSED_FOLDER'] = PROCESSED_FOLDER
app.config['DATASET_FOLDER'] = DATASET_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # Max upload size of 200 MB

# Admin credentials
ADMIN_USERNAME = 'admin'
ADMIN_PASSWORD = 'admin'

app.jinja_env.globals.update(datetime=datetime)
socketio = SocketIO(app)
socketio.init_app(app)  # Pastikan ini diatur setelah inisialisasi socketio

ALLOWED_EXTENSIONS = {'pdf', 'jpg', 'jpeg', 'bmp', 'png', 'docx', 'pptx', 'xlsx', 'txt'}


# Create necessary directories
for folder in [PROCESSED_FOLDER, DATASET_FOLDER, 'logs', 'temporary']:
    os.makedirs(folder, exist_ok=True)
    

EXPECTED_MIME_TYPES = {
    'pdf': ['application/pdf'],
    'jpg': ['image/jpeg'],
    'jpeg': ['image/jpeg'],
    'bmp': ['image/bmp'],
    'png': ['image/png'],
    'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
    'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
    'xlsx': ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
    'txt': ['text/plain']
}

if not os.path.exists('temporary'):
    os.makedirs('temporary')
    
if not os.path.exists('logs'):
    os.makedirs('logs')
#print(f"Logs directory exists: {os.path.exists('logs')}")

# Function to save metadata into a CSV file as a dataset with labels
def save_metadata_to_dataset(filename, metadata, label="normal"):
    dataset_path = os.path.join(DATASET_FOLDER, 'dataset.csv')

    # Define dataset fields, including the label
    fieldnames = [
        'filename', 'file_hash', 'page_hashes', 'metadata', 'watermarked', 'timestamp', 'text_content', 
        'file_size', 'num_pages', 'camera_model', 'camera_make', 'date_time_original', 'hash_md5',
        'hash_sha1', 'hash_sha256', 'hash_sha512', 'last_accessed', 'last_modified', 'created', 'label'
    ]

    # Load existing dataset to check for duplicates
    if os.path.isfile(dataset_path):
        df = pd.read_csv(dataset_path)
        existing_rows = df[(df['filename'] == filename) & (df['file_hash'] == metadata.get('file_hash', ''))]
        
        # Jika data sudah ada, skip penyimpanan
        if not existing_rows.empty:
            print(f"Metadata for {filename} already exists in the dataset. Skipping.")
            return

    # Open the CSV file in append mode, ensuring the header is written only if the file is empty or new
    with open(dataset_path, mode='a', newline='', encoding='utf-8') as csv_file:
        writer = csv.DictWriter(csv_file, fieldnames=fieldnames)

        # Check if the file is empty to write headers
        if os.stat(dataset_path).st_size == 0:
            writer.writeheader()  # Write headers if file is empty

        # Prepare the row with extracted metadata and label
        row = {
            'filename': filename,
            'file_hash': metadata.get('file_hash', ''),
            'page_hashes': json.dumps(metadata.get('page_hashes', {})),
            'metadata': json.dumps(metadata.get('metadata', {})),
            'watermarked': metadata.get('watermarked', False),
            'timestamp': metadata.get('timestamp', ''),
            'text_content': metadata.get('text_content', ''),
            'file_size': metadata.get('FileSizeBytes', 0),
            'num_pages': metadata.get('Pages', 0),
            'camera_model': metadata.get('CameraModel', ''),
            'camera_make': metadata.get('CameraMake', ''),
            'date_time_original': metadata.get('DateTimeOriginal', ''),
            'hash_md5': metadata.get('Hash_MD5', ''),
            'hash_sha1': metadata.get('Hash_SHA1', ''),
            'hash_sha256': metadata.get('Hash_SHA256', ''),
            'hash_sha512': metadata.get('Hash_SHA512', ''),
            'last_accessed': metadata.get('LastAccessed', ''),
            'last_modified': metadata.get('LastModified', ''),
            'created': metadata.get('Created', ''),
            'label': label  # Add the label here
        }
        writer.writerow(row)
    print(f"Metadata for {filename} added to dataset with label '{label}'.")


# Load dataset for training or updating model
def load_dataset():
    dataset_path = os.path.join(DATASET_FOLDER, 'dataset.csv')
    if os.path.isfile(dataset_path):
        df = pd.read_csv(dataset_path)
        
        # Pastikan kolom yang diperlukan ada
 #       required_columns = ['file_size', 'num_pages', 'camera_model', 'camera_make', 'hash_md5', 'hash_sha1', 'hash_sha256', 'hash_sha512']
        required_columns = ['file_size', 'num_pages', 'hash_md5', 'hash_sha1', 'hash_sha256', 'hash_sha512']
        for col in required_columns:
            if col not in df.columns:
                print(f"Warning: Column {col} not found in dataset. Check save_metadata_to_dataset function.")
                return np.array([])  # Mengembalikan array kosong jika kolom tidak ditemukan
        
        # Ubah field non-numerik ke panjang string agar bisa diolah
        features = df[required_columns]
        features = features.applymap(lambda x: len(str(x)) if isinstance(x, str) else x)
        return features.to_numpy()
    return np.array([])


# Train or update the anomaly detection model
def train_anomaly_model():
    data = load_dataset()
    if data.size == 0:
        print("No data available for training.")
        return

    scaler = StandardScaler()
    data_scaled = scaler.fit_transform(data)
    model = IsolationForest(contamination=0.1, random_state=42)
    model.fit(data_scaled)

    joblib.dump(model, MODEL_PATH)
    joblib.dump(scaler, SCALER_PATH)
    print("Anomaly detection model trained and saved.")


# Load trained model for detecting anomalies
def load_anomaly_model():
    if os.path.isfile(MODEL_PATH) and os.path.isfile(SCALER_PATH):
        model = joblib.load(MODEL_PATH)
        scaler = joblib.load(SCALER_PATH)
        return model, scaler
    else:
        print("Model not found. Training a new model.")
#        train_anomaly_model()
        train_anomaly_model_with_known_anomalies()
        return load_anomaly_model()
    

@app.route('/detect-anomalies', methods=['GET'])
def detect_anomalies():
    try:
        data = load_dataset()
        if data.size == 0:
            return render_template('anomaly_results.html', anomaly_log=[])  # Tetap arahkan ke halaman hasil meskipun tidak ada data

        model, scaler = load_anomaly_model()
        data_scaled = scaler.transform(data)
        predictions = model.predict(data_scaled)

        # Identify potentially tampered files (label -1)
        dataset_path = os.path.join(DATASET_FOLDER, 'dataset.csv')
        df = pd.read_csv(dataset_path)
        anomalous_files = df['filename'][predictions == -1].tolist()

        # Log results of anomaly detection if not already present
        logged_filenames = set()
        if os.path.exists(ANOMALY_LOG_PATH):
            with open(ANOMALY_LOG_PATH, 'r', encoding='utf-8') as log_file:
                logged_filenames = {row['filename'] for row in csv.DictReader(log_file)}

        with open(ANOMALY_LOG_PATH, mode='a', newline='', encoding='utf-8') as log_file:
            log_writer = csv.writer(log_file)
            if os.stat(ANOMALY_LOG_PATH).st_size == 0:
                log_writer.writerow(['timestamp', 'filename', 'result'])

            for filename in df['filename']:
                if filename not in logged_filenames:
                    result = "Anomaly Detected" if filename in anomalous_files else "No Anomaly"
                    log_writer.writerow([datetime.now().strftime("%Y-%m-%d %H:%M:%S"), filename, result])

        # Load log data for display in the template, ordered by descending timestamp
        anomaly_log = []
        with open(ANOMALY_LOG_PATH, 'r', encoding='utf-8') as log_file:
            log_reader = csv.DictReader(log_file)
            anomaly_log = sorted(log_reader, key=lambda x: x['timestamp'], reverse=True)  # Default descending order

        return render_template('anomaly_results.html', anomaly_log=anomaly_log)
    except Exception as e:
        flash(f"Error detecting anomalies: {e}", "error")
        return redirect(url_for('dashboard'))


# Route to render the blockchain page
@app.route('/blockchain')
def blockchain_visualization():
    return render_template('blockchain.html', blockchain=blockchain.chain)


# Utility: Save file hash using multiple algorithms (MD5, SHA-1, SHA-256, SHA-512)
def get_file_hash(filepath, algorithm='sha256'):
    hash_algorithms = {
        'md5': hashlib.md5(),
        'sha1': hashlib.sha1(),
        'sha256': hashlib.sha256(),
        'sha512': hashlib.sha512()
    }
    hasher = hash_algorithms.get(algorithm, hashlib.md5())
    
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hasher.update(chunk)
    
    # Ensure the hash is returned as a string (hexadecimal format)
    return hasher.hexdigest()


@app.route('/documentation')
def documentation():
    return render_template('documentation.html')


@app.route('/check-session')
def check_session():
    username = session.get('username', 'No user logged in')
    return f'Current user: {username}'


# Route: Reset confirmation page
@app.route('/reset', methods=['GET', 'POST'])
def reset():
    if 'logged_in' not in session:
        return redirect(url_for('login'))
    
    if request.method == 'POST':
        if request.form.get('confirm_reset') == 'yes':
            try:
                # Clear all files in uploads/ and processed/
                #shutil.rmtree(app.config['UPLOAD_FOLDER'], ignore_errors=True)
                shutil.rmtree(app.config['PROCESSED_FOLDER'], ignore_errors=True)
                #shutil.rmtree(app.config['DATASET_FOLDER'], ignore_errors=True)
                #os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
                #os.makedirs(app.config['DATASET_FOLDER'], exist_ok=True)
                os.makedirs(app.config['PROCESSED_FOLDER'], exist_ok=True)

                # Clear the log file
                with open('logs/file_activity_log.json', 'w') as log_file:
                    log_file.write('')  # Reset the log file
                    
                # Reset the app_error.log file by clearing its content
                with open('logs/app_error.log', 'w') as error_log_file:
                    error_log_file.write('')  # Reset the app_error_log
                
                # Reset the dataset file by clearing its content
                #with open('dataset/dataset.csv', 'w') as dataset_file:
                #    dataset_file.write('')  # Reset the dataset
                    
                #flash('System has been reset successfully!', 'success')
                session.pop('logged_in', None)
                return redirect(url_for('login'))
            except Exception as e:
                flash(f'Error during reset: {e}', 'error')
                return redirect(url_for('reset'))
        else:
            #flash('Reset operation canceled.', 'info')
            return redirect(url_for('dashboard'))

    return render_template('reset.html')


# Handle geolocation failure
def get_geolocation(ip_address):
    try:
        response = requests.get(f"https://ipinfo.io/{ip_address}/json")
        if response.status_code == 200:
            data = response.json()
            return data.get("city"), data.get("region"), data.get("country")
    except Exception as e:
        print(f"Error retrieving geolocation: {str(e)}")
    return None, None, None


def extract_metadata(filepath, file_type):
    metadata = {}

    # Common metadata for all files
    metadata['Filename'] = os.path.basename(filepath)
    metadata['FileSizeBytes'] = os.path.getsize(filepath)  # in bytes
    metadata['MimeType'] = mimetypes.guess_type(filepath)[0]
    metadata['LastAccessed'] = datetime.fromtimestamp(os.path.getatime(filepath)).isoformat()
    metadata['LastModified'] = datetime.fromtimestamp(os.path.getmtime(filepath)).isoformat()
    metadata['Created'] = datetime.fromtimestamp(os.path.getctime(filepath)).isoformat()
    metadata['Hash_MD5'] = hashlib.md5(open(filepath, 'rb').read()).hexdigest()
    metadata['Hash_SHA1'] = hashlib.sha1(open(filepath, 'rb').read()).hexdigest()
    metadata['Hash_SHA256'] = hashlib.sha256(open(filepath, 'rb').read()).hexdigest()
    metadata['Hash_SHA512'] = hashlib.sha512(open(filepath, 'rb').read()).hexdigest()   

    # Specific metadata by file type
    if file_type == 'pdf':
        with open(filepath, 'rb') as f:
            reader = PdfReader(f)
            metadata['IsEncrypted'] = reader.is_encrypted
            metadata['Pages'] = len(reader.pages)
            metadata['Metadata'] = {key: str(value) for key, value in (reader.metadata or {}).items()}
            metadata['Producer'] = reader.metadata.get('/Producer', None)
            metadata['Creator'] = reader.metadata.get('/Creator', None)
            metadata['Author'] = reader.metadata.get('/Author', None)
            metadata['Subject'] = reader.metadata.get('/Subject', None)
            metadata['Title'] = reader.metadata.get('/Title', None)
            metadata['PageContent'] = [{'PageNumber': i + 1, 'Text': (page.extract_text() or '')} for i, page in enumerate(reader.pages)]

    elif file_type in ['jpg', 'jpeg', 'bmp', 'png']:
        with Image.open(filepath) as img:
            metadata['Width'], metadata['Height'] = img.size
            metadata['ColorMode'] = img.mode
            metadata['Format'] = img.format
            metadata['ImageDPI'] = img.info.get('dpi', 'N/A')
            metadata['Transparency'] = img.info.get('transparency', None)
            metadata['Compression'] = img.info.get('compression', None)
            metadata['ICCProfile'] = img.info.get('icc_profile', None)
            
            # Histogram data
            metadata['Histogram'] = img.histogram()

            # Image statistics
            stat = ImageStat.Stat(img)
            metadata['Mean'] = stat.mean  # Mean pixel values
            metadata['Median'] = stat.median  # Median pixel values
            metadata['StandardDeviation'] = stat.stddev  # Standard deviation of pixel values
            metadata['Extrema'] = stat.extrema  # Min and max pixel values per channel

            # EXIF data if available
            exif_data = img._getexif()
            if exif_data:
                exif = {ExifTags.TAGS.get(tag, tag): str(value) for tag, value in exif_data.items()}
                metadata['EXIF'] = exif
                
                # Specific EXIF data
                metadata['CameraMake'] = exif.get('Make', None)
                metadata['CameraModel'] = exif.get('Model', None)
                metadata['ExposureTime'] = exif.get('ExposureTime', None)
                metadata['FNumber'] = exif.get('FNumber', None)
                metadata['ISO'] = exif.get('ISOSpeedRatings', None)
                metadata['DateTimeOriginal'] = exif.get('DateTimeOriginal', None)
                metadata['Flash'] = exif.get('Flash', None)
                metadata['FocalLength'] = exif.get('FocalLength', None)
                metadata['GPSInfo'] = exif.get('GPSInfo', None)  # GPS metadata if available

            # Additional details for image formats that support it
            if file_type in ['png', 'tiff']:
                metadata['Gamma'] = img.info.get('gamma', None)
                metadata['ResolutionUnit'] = exif.get('ResolutionUnit', None) if exif_data else None

            # For images with alpha channels or transparency
            if 'transparency' in img.info:
                metadata['Transparency'] = img.info['transparency']
            if img.mode == 'RGBA':
                metadata['AlphaChannel'] = 'Present'

    elif file_type == 'docx':
        doc = Document(filepath)
        core_properties = doc.core_properties
        metadata['ParagraphCount'] = len(doc.paragraphs)
        metadata['WordCount'] = sum(len(paragraph.text.split()) for paragraph in doc.paragraphs)
        metadata['CharacterCount'] = sum(len(paragraph.text) for paragraph in doc.paragraphs)
        metadata.update({
            'Author': getattr(core_properties, 'author', None),
            'Created': core_properties.created.isoformat() if core_properties.created else None,
            'Identifier': getattr(core_properties, 'identifier', None),
            'Keywords': getattr(core_properties, 'keywords', None),
            'Language': getattr(core_properties, 'language', None),
            'LastModifiedBy': getattr(core_properties, 'last_modified_by', None),
            'LastPrinted': core_properties.last_printed.isoformat() if core_properties.last_printed else None,
            'Modified': core_properties.modified.isoformat() if core_properties.modified else None,
            'Revision': getattr(core_properties, 'revision', None),
            'Subject': getattr(core_properties, 'subject', None),
            'Title': getattr(core_properties, 'title', None),
            'Category': getattr(core_properties, 'category', None)
        })
        metadata['Sections'] = [{'SectionNumber': i + 1, 'Text': paragraph.text} for i, paragraph in enumerate(doc.paragraphs)]

    elif file_type == 'pptx':
        ppt = Presentation(filepath)
        core_properties = ppt.core_properties
        metadata['SlideCount'] = len(ppt.slides)
        metadata.update({
            'Title': core_properties.title,
            'Author': core_properties.author,
            'Created': core_properties.created.isoformat() if core_properties.created else None,
            'LastModified': core_properties.modified.isoformat() if core_properties.modified else None,
            'Subject': core_properties.subject,
            'Keywords': core_properties.keywords,
            'LastPrinted': core_properties.last_printed.isoformat() if core_properties.last_printed else None,
        })
        metadata['Slides'] = [{'SlideNumber': i + 1, 'Text': ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])} for i, slide in enumerate(ppt.slides)]
        
    elif file_type == 'xlsx':
        if file_type == 'xlsx':
            wb = load_workbook(filepath, data_only=True)
            metadata['Sheets'] = wb.sheetnames
            sheet_metadata = []

            for sheet in wb.sheetnames:
                sheet_data = {'SheetName': sheet, 'Rows': [], 'Columns': []}
                ws = wb[sheet]
                sheet_data['MaxRows'] = ws.max_row
                sheet_data['MaxColumns'] = ws.max_column

                # Looping untuk membaca isi data dari beberapa baris/kolom pertama
                data_preview = []
                for row in ws.iter_rows(min_row=1, max_row=min(5, ws.max_row), min_col=1, max_col=min(5, ws.max_column)):
                    row_data = []
                    for cell in row:
                        # Pastikan cell adalah objek Cell yang valid
                        if isinstance(cell, Cell):
                            row_data.append(cell.value)
                        else:
                            row_data.append(None)
                    data_preview.append(row_data)

                sheet_data['DataPreview'] = data_preview
                sheet_metadata.append(sheet_data)

            metadata['SheetDetails'] = sheet_metadata

    elif file_type == 'txt':
        with open(filepath, 'r', encoding='utf-8') as f:
            content = f.read()
            metadata.update({
                'SizeInBytes': os.path.getsize(filepath),
                'WordCount': len(content.split()),
                'LineCount': content.count('\n'),
                'CharacterCount': len(content),
                'Encoding': f.encoding,
                'ContentPreview': content[:500]
            })

    return metadata


# Function: Add Invisible Watermark to PDF
def add_invisible_watermark(input_pdf, output_pdf, watermark_text):
    try:
        pdf = pikepdf.open(input_pdf)
        watermark_info = f'rahasia {watermark_text}'
        pdf.docinfo['/Watermark'] = watermark_info
        pdf.save(output_pdf)
        flash(f"Watermark added to {input_pdf}", 'success')
    except PdfError:
        flash(f"Failed to process {input_pdf}. The file may be corrupted.", 'error')

        
# Function: Verify file integrity
def verify_file_integrity(filepath, original_hash):
    current_hash = get_file_hash(filepath)
    return current_hash == original_hash


@app.route('/upload', methods=['GET', 'POST'])
def upload():
    if 'logged_in' not in session:
        return redirect(url_for('login'))

    if request.method == 'POST':
        ip_address = request.remote_addr
        print(f"IP address: {ip_address}")
        
        # Cek apakah input berupa file tunggal atau folder
        files = []
        if 'file_single' in request.files and request.files['file_single'].filename != '':
            files.append(request.files['file_single'])
        elif 'file_folder[]' in request.files:
            files = request.files.getlist('file_folder[]')

        # Periksa apakah ada file yang diunggah
        if not files:
            flash('No file selected.', 'error')
            return redirect(url_for('upload'))

        watermark_choice = request.form.get('add_watermark')  # Get watermark choice from form

        for file in files:
            if file:
                # Dapatkan nama file dan simpan ke folder temporary untuk diproses
                filename = secure_filename(os.path.basename(file.filename))
                temp_filepath = os.path.join('temporary', filename)
                processed_filepath = os.path.join(app.config['PROCESSED_FOLDER'], filename)

                # Simpan file sementara untuk proses hash
                file.save(temp_filepath)

                if allowed_file(filename, temp_filepath):
                    file_hash = get_file_hash(temp_filepath)
                    
                    # Cek apakah file dengan hash yang sama sudah pernah diproses
                    if is_file_already_processed(file_hash):
                        os.remove(temp_filepath)  # Hapus file sementara
                        flash(f"File '{filename}' has already been processed.", 'warning')
                        continue  # Lanjutkan ke file berikutnya

                    # Ekstraksi metadata
                    metadata = extract_metadata(temp_filepath, filename.rsplit('.', 1)[1].lower())
                    metadata['file_hash'] = file_hash
                    metadata['timestamp'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                    
                    # Jika file adalah PDF, tambahkan page hash dan opsi watermark
                    if filename.endswith('.pdf'):
                        page_hashes = get_pdf_page_hashes(temp_filepath)
                        print("Page hashes during upload:", page_hashes)  # Debug print
                        metadata['page_hashes'] = page_hashes
                        
                        if watermark_choice == 'yes':  # Hanya tambahkan watermark jika dipilih
                            timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                            add_invisible_watermark(temp_filepath, processed_filepath, timestamp)
                            metadata['watermarked'] = True
                        else:
                            shutil.move(temp_filepath, processed_filepath)
                            metadata['watermarked'] = False
                    else:
                        # Pindahkan file ke folder processed jika bukan PDF atau tanpa watermark
                        shutil.move(temp_filepath, processed_filepath)

                    # Simpan metadata dan log aktivitas
                    save_metadata(filename, metadata)
                    save_metadata_to_dataset(filename, metadata)
                    
                    log_file_activity('upload', filename, ip_address)
                    detect_anomalies()

                    flash(f'Metadata for {filename} saved.', 'success')
                else:
                    os.remove(temp_filepath)
                    flash(f'File {filename} is not valid.', 'error')

        return redirect(url_for('upload'))

    return render_template('upload.html')


@app.route('/dashboard')
def dashboard():
    if 'logged_in' not in session:
        return redirect(url_for('login'))
    return render_template('dashboard.html')  # Halaman dengan navigasi menu


# Route: Login Page
@app.route('/', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        if username == ADMIN_USERNAME and password == ADMIN_PASSWORD:
            session['logged_in'] = True
            session['username'] = username  # Store the username in session
            ip_address = request.remote_addr
            log_file_activity('login', 'N/A', ip_address)
            
            print(f"User {username} logged in.")  # Debug log
            return redirect(url_for('upload'))
        else:
            flash('Invalid credentials', 'error')
    return render_template('login.html')


# Function: Validate file existence in processed folder
def validate_file(filename):
    processed_files = os.listdir(app.config['PROCESSED_FOLDER'])
    if filename not in processed_files:
        flash(f"{filename} is not yet processed. Please upload the file first.", 'warning')
        return False
    return True

import magic


# Function: Check allowed file types
def allowed_file(filename, filepath):
    ext = filename.rsplit('.', 1)[1].lower()
    mime = magic.from_file(filepath, mime=True)
    if ext in ALLOWED_EXTENSIONS and mime in EXPECTED_MIME_TYPES[ext]:
        return True
    return False


@app.route('/validate-file', methods=['GET', 'POST'])
def validate_file_action():
    if 'logged_in' not in session:
        return redirect(url_for('login'))

    if request.method == 'POST':
        file = request.files['file']
        ip_address = request.remote_addr
        validation_result = []

        if file:
            filename = secure_filename(file.filename)
            temp_filepath = os.path.join('temporary', filename)

            # Simpan file sementara di temporary folder
            file.save(temp_filepath)

            # Validasi apakah file termasuk tipe yang diizinkan
            if allowed_file(filename, temp_filepath):
                file_hash = get_file_hash(temp_filepath)

                # Dapatkan daftar file yang telah diproses
                processed_files = [f for f in os.listdir(app.config['PROCESSED_FOLDER']) if f.endswith('.json')]

                # Jika file adalah PDF, dapatkan hash per halaman untuk validasi
                if filename.endswith('.pdf'):
                    uploaded_page_hashes = get_pdf_page_hashes(temp_filepath)
                    # Kirim ke halaman yang sama untuk memilih file perbandingan
                    return render_template(
                        'validate.html', 
                        uploaded_file=temp_filepath, 
                        processed_files=processed_files, 
                        uploaded_page_hashes=uploaded_page_hashes
                    )

                # Kirim file yang diunggah dan daftar file yang diproses ke halaman yang sama
                return render_template(
                    'validate.html', 
                    uploaded_file=temp_filepath, 
                    processed_files=processed_files
                )
            else:
                os.remove(temp_filepath)
                flash(f'File {filename} is not valid.', 'error')
                return redirect(url_for('validate_file_action'))

    return render_template('validate.html', validation_result=None)


@app.route('/compare-file', methods=['POST'])
def compare_file():
    if 'logged_in' not in session:
        return redirect(url_for('login'))

    uploaded_file = request.form['uploaded_file']
    selected_file = request.form['selected_file']
    ip_address = request.remote_addr
    validation_result = []
    page_hashes_match = True

    if uploaded_file and selected_file:
        # Validasi file yang dipilih untuk perbandingan
        uploaded_page_hashes = get_pdf_page_hashes(uploaded_file)

        # Buka metadata file yang diproses (file JSON)
        selected_metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], f'{selected_file}.json')
        with open(selected_metadata_path, 'r') as f:
            metadata = json.load(f)
            if 'page_hashes' in metadata:
                for page, hash_val in uploaded_page_hashes.items():
                    if metadata['page_hashes'].get(page) != hash_val:
                        page_hashes_match = False
                        validation_result.append(f"Page {page} hash does not match.")
            else:
                validation_result.append(f"Processed file {selected_file} has no page hashes.")

        # Catat aktivitas validasi di log
        log_file_activity('validate', uploaded_file, ip_address)

        # Hapus file yang diunggah setelah validasi selesai
        os.remove(uploaded_file)

        # Tambahkan hasil validasi ke halaman
        if page_hashes_match:
            validation_result.append("File is valid and page hashes match.")
        else:
            validation_result.append("Validation failed: Page hashes do not match.")

        return render_template('validate.html', validation_result=validation_result)

    flash('No file selected for comparison.', 'error')
    return redirect(url_for('validate_file_action'))


# Route: Tracking log
@app.route('/logs', methods=['GET'])
def logs():
    # Ambil parameter sorting dan page dari query string
    sort_order = request.args.get('sort', 'desc')  # Default to descending
    page = request.args.get(get_page_parameter(), type=int, default=1)  # Default page number
    per_page = 20  # Jumlah log per halaman

    # Load logs dari file
    logs = []
    try:
        with open('logs/file_activity_log.json', 'r') as log_file:
            for line in log_file:
                logs.append(json.loads(line))
        print(f"Loaded logs: {logs}")  # Debug log
    except FileNotFoundError:
        flash('No log file found.', 'error')
    except json.JSONDecodeError as e:
        flash(f"Error reading log file: {e}", 'error')
        print(f"Error reading log file: {e}")  # Debug log

    # Sort logs berdasarkan timestamp
    reverse_order = (sort_order == 'desc')
    logs = sorted(logs, key=lambda x: x['timestamp'], reverse=reverse_order)

    # Implementasi pagination
    total = len(logs)
    start = (page - 1) * per_page
    end = start + per_page
    logs_to_show = logs[start:end]

    # Buat pagination object
    pagination = Pagination(page=page, total=total, per_page=per_page, record_name='logs', css_framework='bootstrap4')

    # Render template dengan logs yang sudah disortir dan dipaginate
    return render_template('logs.html', logs=logs_to_show, pagination=pagination, sort_order=sort_order)


# Route: Validate file
@app.route('/validate/<filename>', methods=['GET'])
def validate(filename):
    # Check if the file exists in the processed folder
    if not validate_file(filename):
        return redirect(url_for('upload'))

    # If file exists, continue with validation
    status = f"Validation successful for {filename}."
    return render_template('validate.html', filename=filename, status=status)


# Route: Logout
@app.route('/logout')
def logout():
    if 'logged_in' in session:
        ip_address = request.remote_addr
        username = session.get('username', 'Unknown')

        # Log activity for logout
        log_file_activity('logout', 'N/A', ip_address)
        
        session.pop('logged_in', None)
        session.pop('username', None)
        
 #   flash('Logged out successfully', 'success')
    return redirect(url_for('login'))


# Protect the routes to ensure the user is logged in
@app.before_request
def require_login():
    allowed_routes = ['login', 'static', 'documentation']  # Allow access to login and static files
    if 'logged_in' not in session and request.endpoint not in allowed_routes:
        return redirect(url_for('login'))

    
# Function: Save metadata to a JSON file
def save_metadata(filename, metadata):
    try:
        metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], f"{filename}.json")
        with open(metadata_path, 'w') as json_file:
            json.dump(metadata, json_file, ensure_ascii=False)
    except Exception as e:
        print(f"Error saving metadata: {e}")
        flash(f"Failed to save metadata for {filename}.", 'error')

        
def file_already_exists(filepath):
    file_hash = get_file_hash(filepath)
    for processed_file in os.listdir(app.config['PROCESSED_FOLDER']):
        if get_file_hash(os.path.join(app.config['PROCESSED_FOLDER'], processed_file)) == file_hash:
            return True
    return False


# Route: Check watermark in files
@app.route('/check-watermark/<filename>')
def check_pdf_watermark(filename):
    filepath = os.path.join(app.config['PROCESSED_FOLDER'], filename)
    watermark = check_watermark(filepath)
    if watermark:
        flash(f"Watermark found in {filename}: {watermark}", 'success')
    else:
        if not watermark:
            flash(f"No watermark found or the file {filename} is corrupted.", 'warning')
    return redirect(url_for('files'))


@app.route('/check-watermark-file', methods=['GET', 'POST'])
def check_watermark_action():
    if 'logged_in' not in session:
        return redirect(url_for('login'))

    if request.method == 'POST':
        file = request.files['file']
        
        ip_address = request.remote_addr
        
        if file:
            filename = secure_filename(file.filename)
            temp_filepath = os.path.join('temporary', filename)
            file.save(temp_filepath)

            # Perbaikan: Tambahkan temp_filepath saat memanggil allowed_file
            if allowed_file(filename, temp_filepath):
                # Check for watermark in the file
                watermark_result = None
                try:
                    watermark_result = check_watermark(temp_filepath)
                    if watermark_result:
                        watermark_result = f"Watermark found: {watermark_result}"
                    else:
                        watermark_result = "No watermark found."
                    
                    # Log activity for watermark check
                    log_file_activity('watermark', filename, ip_address)
                    
                except Exception as e:
                    watermark_result = f"Error checking watermark: {str(e)}"

                # Clean up the temporary file
                os.remove(temp_filepath)

                return render_template('check_watermark.html', watermark_result=watermark_result)
            else:
                os.remove(temp_filepath)
                flash(f'File {filename} is not valid.', 'error')
                return redirect(url_for('check_watermark_action'))

    return render_template('check_watermark.html')


@app.route('/file/<filename>/details', methods=['GET'])
def file_details(filename):
    metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], f"{filename}.json")
    
    if not os.path.exists(metadata_path):
        flash(f"File {filename} metadata not found.", 'error')
        return redirect(url_for('files'))

    # Load metadata
    with open(metadata_path, 'r') as json_file:
        metadata = json.load(json_file)
    
    # Tambahkan konvensional hashes dan blockchain jika belum ada
    file_path = os.path.join(app.config['PROCESSED_FOLDER'], filename)
    if not metadata.get('conventional_hashes'):
        metadata['conventional_hashes'] = generate_hash_report(file_path)
        blockchain_hash = add_file_to_blockchain(metadata['conventional_hashes']['sha256'])
        metadata['blockchain_hash'] = blockchain_hash
        save_metadata(filename, metadata)
    
    # Pastikan page_hashes terdefinisi, meskipun kosong
#    if 'page_hashes' not in metadata:
#        metadata['page_hashes'] = {}

    return render_template('file_details.html', filename=filename, metadata=metadata)


# Route: List processed files
@app.route('/files', methods=['GET'])
def files():
    if 'logged_in' not in session:
        return redirect(url_for('login'))

    sort_order = request.args.get('sort', 'desc')
    reverse_order = (sort_order == 'desc')

    processed_files = [f for f in os.listdir(app.config['PROCESSED_FOLDER']) if f.endswith('.json')]
    file_details = []

    for processed_file in processed_files:
        metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], processed_file)
        
        # Check if the file is not empty and valid JSON
        try:
            with open(metadata_path, 'r') as json_file:
                if os.stat(metadata_path).st_size == 0:
                    logging.error(f"Metadata file {processed_file} is empty.")
                    continue
                metadata = json.load(json_file)
                file_info = {
                    'filename': processed_file.replace('.json', ''),
                    'file_hash': metadata.get('file_hash', 'Unknown'),
                    'page_hashes': metadata.get('page_hashes', {}),
                    'metadata': metadata,
                    'watermarked': metadata.get('watermarked', False)  # Check if watermarked
                }
                file_details.append(file_info)
        except json.JSONDecodeError as e:
            logging.error(f"Error decoding JSON in file {processed_file}: {e}")
            continue  # Skip file if JSON is invalid

    # Sort the files based on the sorting order
    file_details.sort(key=lambda x: x['filename'], reverse=reverse_order)

    return render_template('files.html', files=file_details, sort_order=sort_order)


@app.route('/file/<filename>/metadata', methods=['GET'])
def file_metadata(filename):
    metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], f"{filename}.json")
    
    if not os.path.exists(metadata_path):
        flash(f"Metadata for {filename} not found.", 'error')
        return redirect(url_for('files'))

    # Load metadata
    with open(metadata_path, 'r') as json_file:
        metadata = json.load(json_file)

    return render_template('file_metadata.html', filename=filename, metadata=metadata)


# Function: Get stored hash for integrity check
def get_stored_hash(filename):
    metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], f"{filename}.json")
    if os.path.exists(metadata_path):
        with open(metadata_path, 'r') as json_file:
            metadata = json.load(json_file)
            return metadata.get('file_hash')
    return None


# Route: Download processed file with logging
@app.route('/download/<filename>')
def download(filename):
    if 'logged_in' not in session:
        return redirect(url_for('login'))
    
    user = session.get('username', 'Unknown')
    ip_address = request.remote_addr
    
    # Log the download activity
    log_download_activity(filename, user, ip_address)
    
    return send_from_directory(app.config['PROCESSED_FOLDER'], filename)


# Function: Log download activity
def log_download_activity(filename, user=None, ip_address=None):
    try:
        user = session.get('username', 'Unknown')  # Get the username from session
        log_entry = {
            'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            'action': 'download',
            'filename': filename,
            'user': user,
            'ip_address': ip_address
        }
        with open('logs/file_activity_log.json', 'a') as log_file:
            json.dump(log_entry, log_file)
            log_file.write("\n")
    except Exception as e:
        print(f"Error logging download activity: {e}")
        #flash(f"Error logging download activity: {e}", 'error')


def log_file_activity(action, filename, ip_address=None):
    user = session.get('username', 'Unknown')  # Ambil username dari session
    log_entry = {
        'timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        'action': action,
        #'filename': os.path.basename(filename),
        'filename': filename,
#        'file_hash': hash_value,
        'user': user,  # Simpan user yang login
        'ip_address': ip_address
    }
    try:
        log_file_path = 'logs/file_activity_log.json'
        with open(log_file_path, 'a') as log_file:
            json.dump(log_entry, log_file)
            log_file.write("\n")
        #print("Log entry successfully written.")
    except Exception as e:
        print(f"Error logging file activity: {e}")
        logging.error(f"Error logging file activity: {e}")
        #flash(f"Error logging file activity: {e}", 'error')


def is_file_already_processed(file_hash):
    processed_files = os.listdir(app.config['PROCESSED_FOLDER'])
    for processed_file in processed_files:
        metadata_path = os.path.join(app.config['PROCESSED_FOLDER'], processed_file)
        
        # Hanya proses file JSON dan pastikan file dalam format teks
        if processed_file.endswith('.json'):
            try:
                with open(metadata_path, 'r', encoding='utf-8') as json_file:
                    metadata = json.load(json_file)
                    if metadata.get('file_hash') == file_hash:
                        return True
            except (json.JSONDecodeError, UnicodeDecodeError, FileNotFoundError) as e:
                logging.error(f"Error reading or parsing JSON file {metadata_path}: {e}")
                continue  # Skip invalid or corrupted files
    return False


# Function: Check for Watermark in PDF
def check_watermark(filepath):
    try:
        pdf = pikepdf.open(filepath)
        watermark = pdf.docinfo.get('/Watermark')
        return watermark
    except PdfError:
        flash(f"Failed to check watermark. The file {os.path.basename(filepath)} may be corrupted.", 'error')
        return None

    
# Function: Securely delete a file using multiple passes
def secure_delete(filepath, passes=3):
    try:
        with open(filepath, 'ba+', buffering=0) as f:
            length = f.tell()
            for _ in range(passes):
                f.seek(0)
                f.write(os.urandom(length))
        os.remove(filepath)
        flash(f'Securely deleted {filepath}', 'success')
    except Exception as e:
        flash(f'Error in securely deleting file: {str(e)}', 'error')

        
# Route: Secure delete a file
@app.route('/secure-delete/<filename>')
def secure_delete_route(filename):
    filepath = os.path.join(app.config['PROCESSED_FOLDER'], filename)
    secure_delete(filepath)
    return redirect(url_for('files'))


# Route: Periodic file integrity check
@app.route('/check-file-integrity')
def check_file_integrity():
    files = os.listdir(app.config['PROCESSED_FOLDER'])
    altered_files = []
    
    for file in files:
        if file.endswith('.json'):  # Skip metadata files
            continue
        stored_hash = get_stored_hash(file)
        if stored_hash:
            filepath = os.path.join(app.config['PROCESSED_FOLDER'], file)
            if not verify_file_integrity(filepath, stored_hash):
                altered_files.append(file)
    
    if altered_files:
        flash(f'The following files have been altered: {", ".join(altered_files)}', 'error')
    else:
        flash('No file integrity issues detected.', 'success')
    
    return redirect(url_for('files'))


# Blockchain for tamper-proof logging
class Blockchain:
    def __init__(self):
        self.chain = []
        self.current_transactions = []
        self.new_block(previous_hash='1', proof=100)

    def new_block(self, proof, previous_hash=None):
        block = {
            'index': len(self.chain) + 1,
            'timestamp': time(),
            'transactions': self.current_transactions,
            'proof': proof,
            'previous_hash': previous_hash or self.hash(self.chain[-1])
        }
        self.current_transactions = []
        self.chain.append(block)
        return block

    def new_transaction(self, file_hash):
        self.current_transactions.append({
            'file_hash': file_hash,
            'timestamp': time(),
        })
        return self.last_block['index'] + 1

    @staticmethod
    def hash(block):
        block_string = json.dumps(block, sort_keys=True).encode()
        return hashlib.sha256(block_string).hexdigest()

    @property
    def last_block(self):
        return self.chain[-1]

    
# Initialize blockchain
blockchain = Blockchain()


# Emit blockchain data when a new block is added
def add_file_to_blockchain(file_hash):
    blockchain.new_transaction(file_hash)
    proof = blockchain.hash(blockchain.last_block)
    new_block = blockchain.new_block(proof)
    block_data = {
        'index': new_block['index'],
        'timestamp': new_block['timestamp'],
        'transactions': new_block['transactions'],
        'proof': new_block['proof'],
        'previous_hash': new_block['previous_hash']
    }
    print("Emitting new_block event with data:", block_data)  # Tambahkan log di sini
    socketio.emit('new_block', block_data)  # Emit event to frontend
    return blockchain.hash(new_block)


def generate_hash_report(filepath):
    return {
        'md5': get_file_hash(filepath, 'md5'),
        'sha1': get_file_hash(filepath, 'sha1'),
        'sha256': get_file_hash(filepath, 'sha256'),
        'sha512': get_file_hash(filepath, 'sha512'),
    }

    
# Function: Generate hash for each page of a PDF    
def get_pdf_page_hashes(filepath):
    try:
        reader = PdfReader(filepath)
        page_hashes = {}
        
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text() or ''  # Extract text or leave empty if no text
            
            if not page_text.strip():  # Jika tidak ada teks, coba lakukan OCR
                images = convert_from_path(filepath, first_page=i+1, last_page=i+1)  # Konversi halaman ke gambar
                page_text = ''
                for img in images:
                    page_text += pytesseract.image_to_string(img)  # OCR pada gambar halaman

            # Buat hash dari teks halaman (baik hasil extract_text atau OCR)
            page_hash = hashlib.sha256(page_text.encode('utf-8')).hexdigest()
            page_hashes[f" {i+1}"] = page_hash
        
        print("Generated page hashes:", page_hashes)  # Debug print
        return page_hashes

    except Exception as e:
        print(f"Error generating page hashes: {str(e)}")
        return {}

# Function: Extract additional features for structural anomalies
def extract_additional_features(filepath):
    features = {}
    page_hashes = get_pdf_page_hashes(filepath)
    
    # 1. Page Hash Consistency: Checks for duplicated or unique page hashes
    unique_hashes = len(set(page_hashes.values()))
    features["page_hash_consistency"] = unique_hashes / len(page_hashes) if page_hashes else 1.0
    
    # 2. Metadata Completeness: Checks if key metadata fields are filled
    metadata = extract_metadata(filepath, 'pdf')
    features["metadata_completeness"] = len([k for k, v in metadata.items() if v]) / len(metadata)
    
    # 3. Detect Broken Font Dictionary: Mark as 1 if missing, 0 otherwise
    features["broken_font_dict"] = 1 if 'Font' not in metadata.get('metadata', {}) else 0

    # 4. Check EOF Marker: Validates if the PDF ends correctly
    try:
        with open(filepath, 'rb') as f:
            f.seek(-5, os.SEEK_END)
            eof_marker = f.read().decode('utf-8')
            features["has_eof_marker"] = 1 if eof_marker == '%%EOF' else 0
    except:
        features["has_eof_marker"] = 0
    
    # 5. Embedded Object Count: Checks for embedded elements that could be anomalous
    features["embedded_object_count"] = len(metadata.get("EmbeddedFiles", []))

    # Additional features could be added here...
    
    return features

# Function to load dataset with base and additional features
def load_dataset_with_additional_features():
    base_features = load_dataset()
    additional_features = []
    dataset_path = os.path.join(DATASET_FOLDER, 'dataset.csv')
    for _, row in pd.read_csv(dataset_path).iterrows():
        filepath = os.path.join(PROCESSED_FOLDER, row['filename'])
        if os.path.exists(filepath):
            additional_features.append(extract_additional_features(filepath))
    return np.hstack([base_features, additional_features])

# Enhanced Anomaly Model Training with Known Anomalies
def train_anomaly_model_with_known_anomalies():
    # Load dataset including additional features
    data = load_dataset_with_additional_features()
    if data.size == 0:
        print("No data available for training.")
        return
    
    scaler = StandardScaler()
    data_scaled = scaler.fit_transform(data)
    model = IsolationForest(contamination=0.05, random_state=42)  # Set to 5% to increase sensitivity
    model.fit(data_scaled)

    joblib.dump(model, MODEL_PATH)
    joblib.dump(scaler, SCALER_PATH)
    print("Enhanced anomaly detection model trained and saved.")





if __name__ == '__main__':
    #app.run(debug=True)
    socketio.run(app, debug=True)